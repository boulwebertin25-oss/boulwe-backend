"""
Extraction du texte des documents et découpage en chunks indexables.
Formats supportés : PDF, Word (.docx), Excel (.xlsx), PowerPoint (.pptx), texte brut.
"""

from pypdf import PdfReader
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from io import BytesIO

from app.embeddings import embed_documents
from app.vectorstore import upsert_chunks

CHUNK_SIZE = 800      # caractères par chunk (approximation simple, pas de tokenizer ici)
CHUNK_OVERLAP = 150   # chevauchement pour ne pas couper une idée en deux


def extract_text_from_pdf(file_bytes: bytes) -> str:
    """Extrait tout le texte d'un fichier PDF."""
    reader = PdfReader(BytesIO(file_bytes))
    pages_text = [page.extract_text() or "" for page in reader.pages]
    return "\n".join(pages_text)


def extract_text_from_docx(file_bytes: bytes) -> str:
    """Extrait le texte des paragraphes et tableaux d'un document Word."""
    doc = Document(BytesIO(file_bytes))

    parts = [p.text for p in doc.paragraphs if p.text.strip()]

    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text for cell in row.cells)
            if row_text.strip():
                parts.append(row_text)

    return "\n".join(parts)


def extract_text_from_xlsx(file_bytes: bytes) -> str:
    """Extrait le contenu de toutes les feuilles d'un classeur Excel."""
    workbook = load_workbook(BytesIO(file_bytes), data_only=True)
    parts = []

    for sheet in workbook.worksheets:
        parts.append(f"--- Feuille : {sheet.title} ---")
        for row in sheet.iter_rows(values_only=True):
            cells = [str(v) for v in row if v is not None]
            if cells:
                parts.append(" | ".join(cells))

    return "\n".join(parts)


def extract_text_from_pptx(file_bytes: bytes) -> str:
    """Extrait le texte de toutes les diapositives d'une présentation PowerPoint."""
    presentation = Presentation(BytesIO(file_bytes))
    parts = []

    for i, slide in enumerate(presentation.slides, start=1):
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                slide_text.append(shape.text_frame.text)
        if slide_text:
            parts.append(f"--- Diapositive {i} ---\n" + "\n".join(slide_text))

    return "\n".join(parts)


def split_into_chunks(text: str, chunk_size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> list[str]:
    """Découpe un texte long en chunks avec chevauchement."""
    text = text.strip()
    if not text:
        return []

    chunks = []
    start = 0
    while start < len(text):
        end = start + chunk_size
        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)
        start += chunk_size - overlap

    return chunks


def ingest_document(file_bytes: bytes, filename: str, collection: str = "general") -> int:
    """
    Pipeline complet d'ingestion : extraction -> découpage -> embeddings -> stockage.
    Retourne le nombre de chunks indexés.
    """
    name = filename.lower()

    if name.endswith(".pdf"):
        text = extract_text_from_pdf(file_bytes)
    elif name.endswith(".docx"):
        text = extract_text_from_docx(file_bytes)
    elif name.endswith(".xlsx"):
        text = extract_text_from_xlsx(file_bytes)
    elif name.endswith(".pptx"):
        text = extract_text_from_pptx(file_bytes)
    else:
        # Fichier texte brut (.txt, .md)
        text = file_bytes.decode("utf-8", errors="ignore")

    chunks = split_into_chunks(text)
    if not chunks:
        return 0

    embeddings = embed_documents(chunks)
    count = upsert_chunks(collection, chunks, embeddings, source_document=filename)
    return count
